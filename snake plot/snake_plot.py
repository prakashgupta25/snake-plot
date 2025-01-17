# -*- coding: utf-8 -*-
"""
Created on Fri Jan 17 12:29:37 2025

@author: PrakashGupta
"""

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches

# Step 1: Create dummy data for the Snake Plot
data = {
    "Category": ["A", "B", "C", "D", "E"],
    "Series1": [5, 7, 9, 6, 8],  # Values for the first line
    "Series2": [6, 5, 7, 9, 6]   # Values for the second line
}

# Step 2: Convert to a DataFrame
df = pd.DataFrame(data)

# Step 3: Load data for the plot
categories = df["Category"]  # x-axis categories
series1 = df["Series1"]  # y-axis values for Series 1
series2 = df["Series2"]  # y-axis values for Series 2

# Step 4: Create a PowerPoint presentation
ppt = Presentation()

# Add a slide with a blank layout
slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide layout
title = slide.shapes.title
title = slide.shapes.title
title.text = "Snake Plot Example"

# Add a line chart to the slide
chart_data = CategoryChartData()
chart_data.categories = categories
chart_data.add_series("Series 1", series1)
chart_data.add_series("Series 2", series2)

# Set chart position and dimensions
x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

# Customize the chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # Set legend position

# Save the presentation
output_path = "Snake_Plot_Presentation.pptx"
ppt.save(output_path)

print(f"Snake Plot exported to PowerPoint successfully! File saved at: {output_path}")
